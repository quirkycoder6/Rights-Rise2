import io
import json
import os
import convertapi
import requests
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Inches
from urllib.parse import quote_plus
from dotenv import load_dotenv
from reportlab.pdfgen import canvas
import aspose.slides as slides
from pymongo import MongoClient


dir_path = '.\myapp\static\presentations'
# convertapi.api_secret = '4NPm2b4CQKdRkdKN'
load_dotenv()
API_KEY = os.getenv('PEXELS_API_KEY')


client = MongoClient('mongodb://127.0.0.1:27017/?compressors=disabled&gssapiServiceName=mongodb')
db = client['pptx_database']
collection = db['pptx_collection']
convertapi.api_secret = '4NPm2b4CQKdRkdKN'

# # Write PPTX file to MongoDB
# def write_to_mongodb(file_path):
#     with open(file_path, 'rb') as file:
#         pptx_content = file.read()
#         collection.insert_one({'file_name': file_path, 'content': pptx_content})

# # Read PPTX file from MongoDB and save to local system
# def read_from_mongodb_save_locally(file_name, save_path):
#     pptx_data = collection.find_one({'file_name': file_name})
#     if pptx_data:
#         pptx_content = pptx_data['content']
#         pptx_path = os.path.join(save_path, file_name)  # Specify save path
        
#         # Write pptx content to a temporary file
#         with open(pptx_path, 'wb') as tmp_pptx:
#             tmp_pptx.write(pptx_content)
        
#         return pptx_path
#     else:
#         return None

# # Delete local PPTX file after inserting to MongoDB
# def delete_local_file(file_path):
#     os.remove(file_path)

# # Convert PPTX to PDF
# def convert_pptx_to_pdf(input_path, output_path):
#     convertapi.convert('pdf', {'File': input_path}, from_format='pptx').save_files(output_path)

def parse_response(response):
    slides = response.split('\n\n')
    slides_content = []
    for slide in slides:
        lines = slide.split('\n')
        title_line = lines[0]
        if ': ' in title_line:
            title = title_line.split(': ', 1)[1]  # Extract the title after 'Slide X: '
        else:
            title = title_line
        content_lines = [line for line in lines[1:] if line != 'Content:']  # Skip line if it is 'Content:'
        content = '\n'.join(content_lines)  # Join the lines to form the content
        # Extract the keyword from the line that starts with 'Keyword:'
        keyword_line = [line for line in lines if 'Keyword:' or 'Keywords:' in line][0]
        keyword = keyword_line.split(': ', 1)[1]
        slides_content.append({'title': title, 'content': content, 'keyword': keyword})
    return slides_content


def search_pexels_images(keyword):
    query = quote_plus(keyword.lower())
    print("Query:", query) # Debug
    PEXELS_API_URL = f'https://api.pexels.com/v1/search?query={query}&per_page=1'
    print("URL:", PEXELS_API_URL) # Debug
    headers = {
        'Authorization': API_KEY
    }
    response = requests.get(PEXELS_API_URL, headers=headers)
    print("Response Status Code:", response.status_code) # Debug
    print("Response Content:", response.text) # Debug
    data = json.loads(response.text)
    if 'photos' in data:
        if len(data['photos']) > 0:
            return data['photos'][0]['src']['medium']
    return None


def delete_first_two_slides(presentation):
    slide_ids = [1, 0]
    for slide_id in slide_ids:
        if slide_id < len(presentation.slides):
            xml_slides = presentation.slides._sldIdLst
            slides = list(xml_slides)
            xml_slides.remove(slides[slide_id])


def create_ppt(slides_content, template_choice, presentation_title, presenter_name, insert_image):
    template_path = os.path.join(dir_path, f"{template_choice}.pptx")

    prs = Presentation(template_path)

    title_slide_layout = prs.slide_layouts[0]
    content_slide_layout = prs.slide_layouts[1]

    # add title slide
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    title.text = presentation_title

    #add subtitle
    subtitle = slide.placeholders[1]
    subtitle.text = f"Presented by {presenter_name}"

    if template_choice == 'dark_modern':
        for paragraph in title.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = 'Times New Roman'
                run.font.color.rgb = RGBColor(255, 165, 0)  # RGB for orange color

    elif template_choice == 'bright_modern':
        for paragraph in title.text_frame.paragraphs:
            for run in paragraph.runs:
                run.font.name = 'Arial'
                run.font.color.rgb = RGBColor(255, 20, 147)  # RGB for deep pink color

    # add content slides
    for slide_content in slides_content:
        slide = prs.slides.add_slide(content_slide_layout)

        for placeholder in slide.placeholders:
            if placeholder.placeholder_format.type == 1:  # Title
                placeholder.text = slide_content['title']
                if template_choice == 'dark_modern':
                    for paragraph in placeholder.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.name = 'Times New Roman'
                            run.font.color.rgb = RGBColor(255, 165, 0)  # RGB for orange color
            elif placeholder.placeholder_format.type == 7:  # Content
                placeholder.text = slide_content['content']
                if template_choice == 'dark_modern':
                    for paragraph in placeholder.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.font.name = 'Times New Roman'
                            run.font.color.rgb = RGBColor(255, 255, 255)  # RGB for white color

        if insert_image:
            # fetch image URL from Pixabay based on the slide's title
            image_url = search_pexels_images(slide_content['keyword'])
            print("Image URL:", image_url) #debug
            if image_url is not None:
                # download the image
                image_data = requests.get(image_url).content
                # load image into BytesIO object
                image_stream = io.BytesIO(image_data)
                # add the image at the specified position
                slide_width = Inches(20)
                slide_height = Inches(15)

                image_width = Inches(8)  # width of image
                image_height = Inches(5)  # height of image

                left = slide_width - image_width  # calculate left position
                top = slide_height - image_height - Inches(4)  # calculate top position

                slide.shapes.add_picture(image_stream, left, top, width=image_width, height=image_height)

    # add credits slide
    slide = prs.slides.add_slide(content_slide_layout)
    if template_choice == 'dark_modern':
        for placeholder in slide.placeholders:
            if placeholder.placeholder_format.type == 1:  # Title
                placeholder.text = "Credits"
                for paragraph in placeholder.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = 'Times New Roman'
                        run.font.color.rgb = RGBColor(255, 165, 0)
            elif placeholder.placeholder_format.type == 7:  # Content
                placeholder.text = "Images provided by Pexels: https://www.pexels.com"
                for paragraph in placeholder.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = 'Times New Roman'
                        run.font.color.rgb = RGBColor(255, 255, 255)

    elif template_choice == 'bright_modern':
        for placeholder in slide.placeholders:
            if placeholder.placeholder_format.type == 1:  # Title
                placeholder.text = "Credits"
                for paragraph in placeholder.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = 'Arial'
                        run.font.color.rgb = RGBColor(255, 20, 147)
            elif placeholder.placeholder_format.type == 7:  # Content
                placeholder.text = "Images provided by Pexels: https://www.pexels.com"
                for paragraph in placeholder.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = 'Arial'
                        run.font.color.rgb = RGBColor(0, 0, 0)

    else:
        for placeholder in slide.placeholders:
            if placeholder.placeholder_format.type == 1:  # Title
                placeholder.text = "Credits"
                for paragraph in placeholder.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = 'Arial'
                        run.font.color.rgb = RGBColor(0, 0, 0)
            elif placeholder.placeholder_format.type == 7:  # Content
                placeholder.text = "Images provided by Pexels: https://www.pexels.com"
                for paragraph in placeholder.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = 'Arial'
                        run.font.color.rgb = RGBColor(0, 0, 0)

    # Delete the first two slides after all new slides have been added
    delete_first_two_slides(prs)

    # Save the presentation
    prs.save(os.path.join('.\myapp\generated', 'generated_presentation.pptx'))

# # File paths
# pptx_file_path = '.\myapp\generated\generated_presentation.pptx'
# save_path = '.'
# pdf_output_path = '1.pdf'  # Define the PDF output path

# # Write to MongoDB and delete local file
# write_to_mongodb(pptx_file_path)
# delete_local_file(pptx_file_path)

# # Read from MongoDB and save locally with a new name
# loaded_pptx_path = read_from_mongodb_save_locally(pptx_file_path, save_path)

# if loaded_pptx_path:
#     print(f'Successfully loaded PPTX from MongoDB: {loaded_pptx_path}')
# else:
#     print(f'Failed to load PPTX from MongoDB: {pptx_file_path}')

# # Convert the loaded PPTX to PDF
# convert_pptx_to_pdf(loaded_pptx_path, pdf_output_path)

# # Print success message if the PDF was created
# if os.path.exists(pdf_output_path):
#     print(f'Successfully converted PPTX to PDF: {pdf_output_path}')

#     # Replace the previous PDF with the new one
#     os.replace(pdf_output_path, 'previous_presentation.pdf')  # Replace the previous PDF
# else:
#     print(f'Failed to convert PPTX to PDF')







#pres = slides.Presentation(".\myapp\generated\generated_presentation.pptx")

# # Convert PPTX to PDF
# #pres.save("national.pdf", slides.export.SaveFormat.PDF)
# def convert_pptx_to_pdf(input_pptx, output_pdf):
#     prs = Presentation(input_pptx)
#     pdf_canvas = canvas.Canvas(output_pdf)

#     for slide in prs.slides:
#         pdf_canvas.setFont("Helvetica", 12)
#         pdf_canvas.drawString(72, 720, f"Slide {prs.slides.index(slide) + 1}")
#         pdf_canvas.showPage()

#     pdf_canvas.save()

# input_pptx = '.\myapp\generated\generated_presentation.pptx'
# output_pdf = 'ppt.pdf'

# convert_pptx_to_pdf(input_pptx, output_pdf)
# def convert():
#     convertapi.convert('pdf', {
#         'File': './myapp/generated/generated_presentation.pptx'
#     }, from_format = 'pptx').save_files('.')


# def writeOne():
#     # Connect to MongoDB
#     client = pymongo.MongoClient("mongodb://127.0.0.1:27017/?compressors=disabled&gssapiServiceName=mongodb")
#     db = client["pptDB"]
#     collection = db["pptx_collection"]

#     # Read the generated PPTX file as binary data
#     with open('.\myapp\generated\generated_presentation.pptx', 'rb') as file:
#         pptx_binary = file.read()

#     # Define metadata for the PPTX file (e.g., title, author, etc.)
#     pptx_metadata = {
#         "title": "Generated Presentation",
#         "author": "Your Name",
#         # Add any other metadata fields as needed
#     }

#     # Create a document to store the PPTX binary and its metadata
#     pptx_document = {
#         "metadata": pptx_metadata,
#         "file": pptx_binary
#     }

#     # Insert the document into the collection
#     collection.insert_one(pptx_document)

#     # Close the MongoDB connection
#     client.close()